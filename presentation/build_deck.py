"""
Builds a 3-slide PowerPoint deck introducing the NALI Migraine Log
("Headway") app to a clinical audience (fellow physicians).

Audience: neurologists, primary-care colleagues, and women's-health
specialists.

Goal: ~5 minute talk -- what the app is, what data it captures, how
it preserves privacy, and how it can support clinical decision-making.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Brand palette -- pulled from the app's AppIcon (Steel Blue) so the deck
# matches the in-app visual identity.
# ---------------------------------------------------------------------------
BRAND_BLUE = RGBColor(0x44, 0x82, 0xB4)        # Steel Blue
BRAND_BLUE_DARK = RGBColor(0x2E, 0x5A, 0x82)
BRAND_BLUE_LIGHT = RGBColor(0xE3, 0xEE, 0xF7)
ACCENT_TEAL = RGBColor(0x2A, 0x9D, 0x8F)
INK = RGBColor(0x1B, 0x2A, 0x3A)               # near-black text
SUBINK = RGBColor(0x4A, 0x5A, 0x6A)            # muted body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HAIRLINE = RGBColor(0xCF, 0xD8, 0xDC)
CARD_BG = RGBColor(0xF6, 0xF9, 0xFC)


# ---------------------------------------------------------------------------
# Presentation skeleton -- 16:9
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]  # fully blank layout -- we draw everything


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if not shadow:
        # python-pptx exposes no clean shadow API; strip the default
        sp = shape.shadow
        sp.inherit = False
    return shape


def add_round(slide, x, y, w, h, fill, line=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    sp = shape.shadow
    sp.inherit = False
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name="Calibri",
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items,
    *,
    size=13,
    color=INK,
    bullet_color=BRAND_BLUE,
    line_spacing=1.25,
    bullet_char="\u25CF",  # filled circle
):
    """Render a list of (lead, body) tuples or plain strings as bullets."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        bullet_run = p.add_run()
        bullet_run.text = f"{bullet_char}  "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color

        if isinstance(item, tuple):
            lead, body = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = (" " if not lead.endswith(" ") else "") + body
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_brand_band(slide, height=Inches(0.18)):
    """Thin steel-blue accent strip along the top edge."""
    add_rect(slide, 0, 0, SLIDE_W, height, BRAND_BLUE)


def add_footer(slide, page_no, total=3):
    add_text(
        slide,
        Inches(0.5),
        SLIDE_H - Inches(0.42),
        Inches(8),
        Inches(0.3),
        "NALI Migraine Log  \u2022  Privacy-first migraine tracking for iPhone, "
        "Apple Watch, and Mac",
        size=9,
        color=SUBINK,
    )
    add_text(
        slide,
        SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(0.42),
        Inches(0.7),
        Inches(0.3),
        f"{page_no} / {total}",
        size=9,
        color=SUBINK,
        align=PP_ALIGN.RIGHT,
    )


def add_brain_logo(slide, cx, cy, size_in=1.4):
    """A simple stylised brain mark drawn from primitives. The real app
    icon is a malformed SVG that cairosvg can't parse, so we render a
    visually similar emblem here instead."""
    s = Inches(size_in)
    x = cx - s / 2
    y = cy - s / 2

    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, s, s)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_BLUE
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Two overlapping ellipses to suggest hemispheres.
    pad = Inches(size_in * 0.18)
    half_w = (s - pad * 2) / 2
    half_h = s - pad * 2
    left = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + pad,
        y + pad,
        half_w,
        half_h,
    )
    left.fill.background()
    left.line.color.rgb = WHITE
    left.line.width = Pt(2.5)
    left.shadow.inherit = False

    right = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + pad + half_w,
        y + pad,
        half_w,
        half_h,
    )
    right.fill.background()
    right.line.color.rgb = WHITE
    right.line.width = Pt(2.5)
    right.shadow.inherit = False

    # Central seam
    seam = slide.shapes.add_connector(
        1,
        x + pad + half_w,
        y + pad + Inches(0.05),
        x + pad + half_w,
        y + pad + half_h - Inches(0.05),
    )
    seam.line.color.rgb = WHITE
    seam.line.width = Pt(2.5)


# ---------------------------------------------------------------------------
# SLIDE 1 -- Title / What it is
# ---------------------------------------------------------------------------
def build_slide_1():
    s = prs.slides.add_slide(BLANK)

    # Steel-blue half panel on the left
    add_rect(s, 0, 0, Inches(5.2), SLIDE_H, BRAND_BLUE)

    # Logo / app mark on the panel
    add_brain_logo(s, Inches(2.6), Inches(2.55), size_in=1.9)

    # App title block on the panel
    add_text(
        s,
        Inches(0.5),
        Inches(3.95),
        Inches(4.4),
        Inches(0.55),
        "NALI Migraine Log",
        size=34,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font_name="Calibri",
    )
    add_text(
        s,
        Inches(0.5),
        Inches(4.55),
        Inches(4.4),
        Inches(0.4),
        "(\u201cHeadway\u201d)",
        size=16,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(5.05),
        Inches(4.4),
        Inches(0.45),
        "iPhone  \u2022  Apple Watch  \u2022  Mac",
        size=14,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Right side -- pitch
    add_text(
        s,
        Inches(5.7),
        Inches(0.7),
        Inches(7.2),
        Inches(0.5),
        "FOR FELLOW CLINICIANS",
        size=11,
        bold=True,
        color=BRAND_BLUE,
    )
    add_text(
        s,
        Inches(5.7),
        Inches(1.05),
        Inches(7.4),
        Inches(1.4),
        "A privacy-first migraine\ntracker that turns patient\ndiaries into clinical signal.",
        size=32,
        bold=True,
        color=INK,
        line_spacing=1.05,
    )

    # Body paragraph
    add_text(
        s,
        Inches(5.7),
        Inches(3.05),
        Inches(7.3),
        Inches(1.6),
        "Patients log headaches in seconds across iPhone, Apple Watch, "
        "Mac, or by voice with Siri. Their entries stay on-device and "
        "sync only through their personal iCloud \u2014 there are no "
        "developer servers, no analytics SDKs, and no third-party data "
        "sharing.",
        size=14,
        color=SUBINK,
        line_spacing=1.35,
    )

    # Three quick stat chips
    chips = [
        ("63+", "automated tests covering the\nprediction & data layer"),
        ("0", "third-party SDKs, trackers,\nor developer-side servers"),
        ("3", "Apple platforms supported\nfrom one shared codebase"),
    ]
    chip_y = Inches(5.0)
    chip_w = Inches(2.3)
    chip_h = Inches(1.6)
    chip_gap = Inches(0.18)
    chip_x = Inches(5.7)
    for i, (big, small) in enumerate(chips):
        x = chip_x + (chip_w + chip_gap) * i
        add_round(s, x, chip_y, chip_w, chip_h, CARD_BG, line=HAIRLINE, radius=0.12)
        add_text(
            s,
            x + Inches(0.1),
            chip_y + Inches(0.15),
            chip_w - Inches(0.2),
            Inches(0.7),
            big,
            size=32,
            bold=True,
            color=BRAND_BLUE_DARK,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            x + Inches(0.15),
            chip_y + Inches(0.85),
            chip_w - Inches(0.3),
            Inches(0.7),
            small,
            size=11,
            color=SUBINK,
            align=PP_ALIGN.CENTER,
            line_spacing=1.2,
        )

    add_footer(s, 1)


# ---------------------------------------------------------------------------
# SLIDE 2 -- What it captures & how it predicts
# ---------------------------------------------------------------------------
def build_slide_2():
    s = prs.slides.add_slide(BLANK)
    add_brand_band(s)

    # Title
    add_text(
        s,
        Inches(0.5),
        Inches(0.35),
        Inches(12),
        Inches(0.55),
        "What it captures \u2014 and what it tells you",
        size=28,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(0.95),
        Inches(12),
        Inches(0.4),
        "Structured data + on-device intelligence, designed around how we actually triage migraine.",
        size=14,
        color=SUBINK,
    )

    # Two-column card layout: left = clinical data captured, right = analytics & prediction
    left_x = Inches(0.5)
    right_x = Inches(6.95)
    col_w = Inches(5.85)
    col_y = Inches(1.6)
    col_h = Inches(5.3)

    # ---------- LEFT CARD ----------
    add_round(s, left_x, col_y, col_w, col_h, CARD_BG, line=HAIRLINE, radius=0.04)
    # Header bar -- flat rectangle sits flush on top of the card; the
    # tiny card radius (0.04) hides the corners cleanly enough.
    add_rect(
        s,
        left_x,
        col_y,
        col_w,
        Inches(0.55),
        BRAND_BLUE,
    )
    add_text(
        s,
        left_x + Inches(0.25),
        col_y + Inches(0.1),
        col_w - Inches(0.5),
        Inches(0.4),
        "Clinical data captured",
        size=16,
        bold=True,
        color=WHITE,
    )

    add_bullets(
        s,
        left_x + Inches(0.3),
        col_y + Inches(0.75),
        col_w - Inches(0.6),
        col_h - Inches(0.85),
        [
            ("Headache events:", "onset/end time, 1\u201310 pain, duration, location, aura."),
            ("Triggers:", "stress, sleep, dehydration, hormones, weather, alcohol, "
                          "caffeine, food, screens, exercise."),
            ("Medications:", "abortives (triptans, NSAIDs, antiemetics, gepants) "
                              "with dose timing for MOH signal."),
            ("Life impact:", "work / school / events missed \u2014 quantifies "
                              "disability burden between visits."),
            ("Apple Health (opt-in):", "sleep, HRV, resting HR, steps, and "
                                       "menstrual flow for context-aware analytics."),
            ("Weather (Open-Meteo):", "barometric pressure, pressure change, "
                                       "humidity, temperature \u2014 attached to every entry."),
            ("Voice + Watch logging:", "\u201cHey Siri, log a migraine\u201d "
                                       "or one tap on the wrist."),
        ],
        size=12,
        color=INK,
        line_spacing=1.25,
    )

    # ---------- RIGHT CARD ----------
    add_round(s, right_x, col_y, col_w, col_h, CARD_BG, line=HAIRLINE, radius=0.04)
    add_rect(
        s,
        right_x,
        col_y,
        col_w,
        Inches(0.55),
        ACCENT_TEAL,
    )
    add_text(
        s,
        right_x + Inches(0.25),
        col_y + Inches(0.1),
        col_w - Inches(0.5),
        Inches(0.4),
        "Analytics & risk prediction",
        size=16,
        bold=True,
        color=WHITE,
    )

    add_bullets(
        s,
        right_x + Inches(0.3),
        col_y + Inches(0.75),
        col_w - Inches(0.6),
        col_h - Inches(0.85),
        [
            ("Overview dashboard:", "frequency, severity buckets (mild/moderate/severe/extreme), "
                                    "migraine-free streak, top trigger, abortives used."),
            ("Severity heatmap:", "60\u201390 day calendar view \u2014 a glance "
                                   "at clustering, weekend patterns, and trends."),
            ("Auto-generated insights:", "only surfaces cards when the underlying "
                                          "signal is statistically meaningful."),
            ("HealthKit correlations:", "sleep on migraine eves vs. controls, HRV in the "
                                         "72-hour prodromal window, and menstrual-cycle phase "
                                         "(highlights the perimenstrual window for estrogen-withdrawal migraine)."),
            ("24-hour risk forecast:", "hybrid scoring \u2014 rule-based from day one, "
                                        "transitions to an on-device CoreML model after "
                                        "\u226520 logged entries."),
            ("Forecast-risk alerts:", "an opt-in push fires \u22482 hours before a "
                                       "high-risk window so patients can pre-empt with their abortive plan."),
        ],
        size=12,
        color=INK,
        line_spacing=1.25,
    )

    add_footer(s, 2)


# ---------------------------------------------------------------------------
# SLIDE 3 -- Privacy posture & clinical relevance
# ---------------------------------------------------------------------------
def build_slide_3():
    s = prs.slides.add_slide(BLANK)
    add_brand_band(s)

    add_text(
        s,
        Inches(0.5),
        Inches(0.35),
        Inches(12),
        Inches(0.55),
        "Privacy by design \u2014 and why it matters at the bedside",
        size=28,
        bold=True,
        color=INK,
    )
    add_text(
        s,
        Inches(0.5),
        Inches(0.95),
        Inches(12),
        Inches(0.4),
        "Patients can hand you better data with zero compromise on confidentiality.",
        size=14,
        color=SUBINK,
    )

    # Top row -- privacy pillars (3 small cards)
    pill_y = Inches(1.55)
    pill_h = Inches(1.55)
    pill_w = Inches(4.07)
    pill_gap = Inches(0.16)
    pill_x0 = Inches(0.5)

    pillars = [
        (
            "On device + iCloud only",
            "All entries persist locally with Core Data and sync exclusively "
            "through the patient\u2019s personal iCloud account. No developer "
            "servers ever see the data.",
        ),
        (
            "No tracking, no SDKs",
            "Zero third-party analytics, advertising, or crash-reporting "
            "libraries. The dependency graph is intentionally empty.",
        ),
        (
            "Apple Health, two-way",
            "Migraines write back to Apple Health as Headache samples "
            "(idempotent, deduplicated) so they appear in the patient\u2019s "
            "own Health record alongside sleep and cycle data.",
        ),
    ]
    for i, (title, body) in enumerate(pillars):
        x = pill_x0 + (pill_w + pill_gap) * i
        add_round(s, x, pill_y, pill_w, pill_h, WHITE, line=HAIRLINE, radius=0.08)
        # Accent stripe
        add_rect(s, x, pill_y, Inches(0.12), pill_h, BRAND_BLUE)
        add_text(
            s,
            x + Inches(0.3),
            pill_y + Inches(0.18),
            pill_w - Inches(0.4),
            Inches(0.4),
            title,
            size=14,
            bold=True,
            color=BRAND_BLUE_DARK,
        )
        add_text(
            s,
            x + Inches(0.3),
            pill_y + Inches(0.6),
            pill_w - Inches(0.4),
            pill_h - Inches(0.65),
            body,
            size=11,
            color=SUBINK,
            line_spacing=1.3,
        )

    # Bottom -- clinical use cases
    use_y = Inches(3.35)
    use_h = Inches(3.55)

    # Background card
    add_round(
        s,
        Inches(0.5),
        use_y,
        Inches(12.33),
        use_h,
        BRAND_BLUE_LIGHT,
        line=HAIRLINE,
        radius=0.04,
    )
    add_text(
        s,
        Inches(0.8),
        use_y + Inches(0.2),
        Inches(11.5),
        Inches(0.45),
        "Where it helps in practice",
        size=18,
        bold=True,
        color=BRAND_BLUE_DARK,
    )

    # Two-column bullet list inside the card
    use_items_left = [
        ("Better history at the visit:", "patients arrive with a structured "
         "diary instead of a vague \u201cmaybe twice a month.\u201d"),
        ("Trigger identification:", "trigger-share percentages and per-trigger "
         "drill-downs make pattern recognition objective."),
        ("Medication-overuse signal:", "triptan/NSAID use in the last 7 days "
         "is tracked explicitly \u2014 catches MOH earlier."),
    ]
    use_items_right = [
        ("Hormonal migraine workup:", "perimenstrual clustering surfaces "
         "automatically when patients log flow in Apple Health."),
        ("Prodromal awareness:", "HRV and sleep deltas in the 72 h pre-onset "
         "window can prompt earlier abortive use."),
        ("Treatment response:", "severity-bucket trends and migraine-free "
         "streak quantify whether a new prophylactic is working."),
    ]

    add_bullets(
        s,
        Inches(0.8),
        use_y + Inches(0.85),
        Inches(5.75),
        use_h - Inches(1),
        use_items_left,
        size=12,
        color=INK,
        bullet_color=BRAND_BLUE_DARK,
        line_spacing=1.3,
    )
    add_bullets(
        s,
        Inches(6.85),
        use_y + Inches(0.85),
        Inches(5.75),
        use_h - Inches(1),
        use_items_right,
        size=12,
        color=INK,
        bullet_color=BRAND_BLUE_DARK,
        line_spacing=1.3,
    )

    # Closing tagline strip
    tag_y = use_y + use_h - Inches(0.55)
    add_round(
        s,
        Inches(0.8),
        tag_y,
        Inches(11.7),
        Inches(0.42),
        BRAND_BLUE_DARK,
        radius=0.4,
    )
    add_text(
        s,
        Inches(0.8),
        tag_y + Inches(0.05),
        Inches(11.7),
        Inches(0.32),
        "Free for patients  \u2022  No account required  \u2022  Works offline  "
        "\u2022  Available on iPhone, Apple Watch, and Mac",
        size=11,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_footer(s, 3)


# ---------------------------------------------------------------------------
build_slide_1()
build_slide_2()
build_slide_3()

OUT = "/workspace/presentation/NALI_Migraine_Log_Clinical_Overview.pptx"
prs.save(OUT)
print(f"wrote {OUT}")
